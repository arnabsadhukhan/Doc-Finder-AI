import os
import io
import re
import time
import string
import pickle
import subprocess
from collections import Counter 

if not os.path.isdir('status/'):
    print('INSTALLING libraries --this is for frst time only')
    import pkg_resources
    import sys
    pip = "pip3" if sys.platform=='linux' else "pip"
    for package in ['pyqt5==5.14.2','pyqtwebengine==5.14.0','pdfminer==20191125','numpy','pandas','python-docx==0.8.10','python-pptx==0.6.18','xlrd==1.2.0','pyenchant==3.0.1','gensim==3.8.2','sklearn']:
        try:
            dist = pkg_resources.get_distribution(package)
            print('{} ({}) is installed'.format(dist.key, dist.version))
        except pkg_resources.DistributionNotFound:
            print('{} is NOT installed'.format(package))
            print('INSTALLING ',package)
            subprocess.call([pip, 'install', package])

import enchant
import numpy as np
import pandas as pd

from PyQt5 import QtCore, QtGui, QtWidgets
from PyQt5.QtCore import QTimer,QBasicTimer,QSize

from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
from pdfminer.converter import TextConverter
from pdfminer.layout import LAParams
from pdfminer.pdfpage import PDFPage

from docx import Document
from pptx import Presentation
from gensim.summarization.summarizer import summarize






stop_words = ['i','tag','is', 'me', 'my', 'myself', 'we', 'our', 'ours', 'ourselves', 'you', 'your', 'yours', 'yourself', 'yourselves', 'he', 'him', 'his', 'himself', 'she', 'her', 'hers', 'herself', 'it', 'its', 'itself', 'they', 'them', 'their', 'theirs', 'themselves', 'what', 'which', 'who', 'whom', 'this', 'that', 'these', 'those', 'am', 'are', 'was', 'were', 'be', 'been', 'being', 'have', 'has', 'had', 'having', 'do', 'does', 'did', 'doing', 'a', 'an', 'the', 'and', 'but', 'if', 'or', 'because', 'as', 'until', 'while', 'of', 'at', 'by', 'for', 'with', 'about', 'against', 'between', 'into', 'through', 'during', 'before', 'after', 'above', 'below', 'to', 'from', 'up', 'down', 'in', 'out', 'on', 'off', 'over', 'under', 'again', 'further', 'then', 'once', 'here', 'there', 'when', 'where', 'why', 'how', 'all', 'any', 'both', 'each', 'few', 'more', 'most', 'other', 'some', 'such', 'no', 'nor', 'not', 'only', 'own', 'same', 'so', 'than', 'too', 'very', 's', 't', 'can', 'will', 'just', 'don', 'should', 'now']



class Ui_SearchWindow(object):
    def __init__(self):
        with open('status/path.txt','r') as f:
                self.datadir = str(f.read())
        f.close()
        self.d = enchant.Dict("en_US")
        self.table = pd.read_excel('status/table.xlsx',index=False)
        for i in self.table.index:
                if type(self.table.loc[i,'keys'])==str:
                    self.table.at[i,'keys'] = eval(self.table.loc[i,'keys'])
                    try:
                    	self.table.at[i,'hash_map'] = eval(self.table.loc[i,'hash_map'])
                    except:
                    	pass

                    self.table.at[i,'text'] = str(self.table.loc[i,'text'])
        self.table = self.table[['file_name','keys','hash_map','text','tag','location']]
    
    def get_tag(self):
        try:
            i =self.table[self.table.file_name==self.showresultslistWidget.currentItem().text()].index[0]
            #print(self.table.loc[i,'file_name'],self.table.loc[i,'tag'])
            index = self.comboBox.findText(self.table.loc[i,'tag'], QtCore.Qt.MatchFixedString)
            if index >= 0:
                self.comboBox.setCurrentIndex(index)
        except:
            print('get tag function error')
    def on_combobox_changed(self,value):
        try:
            i =self.table[self.table.file_name==self.showresultslistWidget.currentItem().text()].index[0]
            if self.table.at[i,'tag'] != value: 
                self.table.at[i,'tag'] = value
                print('changed to',value) 
                self.table.to_excel('status/table.xlsx',index=False)

        except:
            print('on_combobox_changed function error')
        
    def preview(self):
        #print("preview")
        #print("file is :" ,self.sort.Priority[self.showresultslistWidget.currentRow()])
        try:
            i =self.table[self.table.file_name==self.showresultslistWidget.currentItem().text()].index[0]
            text = str(self.table.loc[i,'text'])
        except:
        	return
        
        seatch_text = re.sub(r'['+re.escape(string.punctuation)+']', ' ',text).lower()
        search_query = str(self.searchquerytextEdit.text()).lower().split(" ")
        i1=[]
        for i, data in enumerate(search_query):
            if len(data)<1 or data in stop_words:
                i1.append(i)
        for i in reversed(i1):
        	del search_query[i]
        self.previewlistView.clear()
        self.previewlistView.appendPlainText("Preview:\n")
        for each_query_word in search_query:
            index=0
            for i in range(2):
                try:
	                index = seatch_text.index(' '+each_query_word+' ',index)+1
	                find_stop = text[index-1:].index('. ')
	                try:
	                    find_space = text[index-1:].index('   ')
	                    print(find_space)
	                    find_stop = min(find_stop,find_space)
	                except:
	                	find_stop=find_stop
	                temp_text = ' '.join(text[index-1:index+find_stop+1].split(" ")[1:-1])+'\n'
	                self.previewlistView.appendPlainText(temp_text)
                except:
                	 pass



    def Summary(self):
        #print("preview")
        #print("file is :" ,self.sort.Priority[self.showresultslistWidget.currentRow()])
        try:
            i =self.table[self.table.file_name==self.showresultslistWidget.currentItem().text()].index[0]
            text = self.table.loc[i,'text']
        except:
        	return
        
        
        self.previewlistView.clear()
        self.previewlistView.appendPlainText("Summary:\n")
        try:
            dis_text = summarize(summarize(text)).split(" ")
            for i,d in enumerate(dis_text):
        	    if len(d)>1:
        		    break
            dis_text = " ".join(dis_text[i:])

            self.previewlistView.appendPlainText(dis_text)
        except:
        	pass
    

    def open_a_file(self):
        #print("open a file")
        #if len(self.searchquerytextEdit.text())>0:
        try:
            location = self.table.index[self.table.file_name==self.showresultslistWidget.currentItem().text()].values
            location = self.table.loc[location,'location'].values[0]
            os.startfile(location)
            self.previewlistView.clear()
            self.previewlistView.appendPlainText('File Location: \n'+ location)
        except:
            opener ="open" if sys.platform == "darwin" else "xdg-open"
            location = self.table.index[self.table.file_name==self.showresultslistWidget.currentItem().text()].values
            
            location = self.table.loc[location,'location'].values[0]
            subprocess.call([opener, location ])
            self.previewlistView.clear()
            self.previewlistView.appendPlainText('File Location: \n'+ location)
            #except:
            #    pass
                    
            
    def search(self):
        #print("searching")
        self.status_directory = 'status/' 
        search_query = str(self.searchquerytextEdit.text())   # change search query

        search_query = search_query.lower().split(" ")
        if len(search_query)>=3 and search_query[-2]== 'tag':
        	search_query = search_query[:-2]
        for i, data in enumerate(search_query):
            if len(data)<=1 or data in stop_words:
                del search_query[i]
        search_keys = [query(quer) for quer in search_query]
        #print(search_keys)

        sort = {'Priority':[],'Times_present':[],'details':[]}

        #print('your query is :',search_query,'\n')
        start = time.time()
        table = self.table
        for file in self.table.index:
            for i,search_key in enumerate(search_keys): 
                if search_key in table.loc[file,'keys']:
                    try:
	                    if search_query[i] in table.loc[file,'hash_map'][search_key].keys():
	                        if table.loc[file,'file_name'] not in list(sort['Priority']):
	                            sort['Priority'].append(table.loc[file,'file_name'])
	                            sort['Times_present'].append(table.loc[file,'hash_map'][search_key][search_query[i]]) 
	                            sort['details'].append([{search_query[i]:table.loc[file,'hash_map'][search_key][search_query[i]]}])
	                            #print('\nsearch found in file :',table.loc[file,'file_name'])
	                            #print('word present:',search_query[i],'   Times: ',table.loc[file,'hash_map'][search_key][search_query[i]])
	                        else:
	                            index  = np.where(table.loc[file,'file_name']==np.array(sort['Priority']))[0][0]
	                            #print(index)
	                            #sort['Priority'].append(table.loc[file,'file_name'])
	                            sort['Times_present'][index] += table.loc[file,'hash_map'][search_key][search_query[i]]
	                            sort['details'][index]=(sort['details'][index]+[{search_query[i]:table.loc[file,'hash_map'][search_key][search_query[i]]}])
	                            #print('\nsearch found in file :',table.loc[file,'file_name'])
	                            #print('word present:',search_query[i],'   Times: ',table.loc[file,'hash_map'][search_key][search_query[i]])
                    except:
                    	pass
        #print( '\n\nTIME TAKEN: ',(time.time()-start)*1000,"ms")
        sort = pd.DataFrame(sort)

        sort = sort.sort_values(by=['Times_present'],ascending=False)

        sort['times_prob'] = sort.Times_present/np.max(sort.Times_present)

        sort['lens_prob'] =[len(each) for each in sort.details]
        sort['lens_prob']/=np.max(sort['lens_prob'])

        sort['order_prob'] =[order_calc(search_query,collect(each)) for each in sort.details]
        sort['order_prob']/=np.max(sort['order_prob'])


        sort['Probability'] = (sort['lens_prob']+sort['times_prob']+sort['order_prob'])/3
        sort = sort.sort_values(by=['Probability'],ascending=False)
        self.showresultslistWidget.clear()
        sort = sort.reset_index()
        sort = sort[['Priority','Probability']]
        self.sort = sort

        search_query = str(self.searchquerytextEdit.text()).lower().split(" ")
        if len(search_query)>=3 and  search_query[-2]=='tag' and search_query[-1] in ["technology","arts","biology","chemistry","economy","literature","physics","sports","others"]:
         
            for i in sort.index:
                
                i_index = self.table.index[self.table.file_name == self.sort.loc[i,'Priority']].tolist()[0]
                
                if self.table.loc[ i_index,'tag'].lower()== search_query[-1]:
                   
                    self.showresultslistWidget.addItem(self.sort.loc[i,'Priority']) 
                else:
                    self.sort.drop(i,axis=0) 
        else:
            for i in self.sort.index:
                self.showresultslistWidget.addItem(self.sort.loc[i,'Priority']) 

        if len(search_query)==2 and  search_query[-2]=='tag' and search_query[-1] in ["technology","arts","biology","chemistry","economy","literature","physics","sports","others"]:
          
            self.showresultslistWidget.clear()
            for i in self.table.index:
            	if self.table.loc[i,'tag'].lower()==search_query[-1]:
                    self.showresultslistWidget.addItem(self.table.loc[i,'file_name']) 

    def setupUi(self, MainWindow):
        MainWindow.setObjectName("MainWindow")
        MainWindow.resize(457, 646)
        self.centralwidget = QtWidgets.QWidget(MainWindow)
        self.centralwidget.setObjectName("centralwidget")
        self.frame = QtWidgets.QFrame(self.centralwidget)
        self.frame.setGeometry(QtCore.QRect(0, 0, 481, 631))
        self.frame.setStyleSheet("QFrame{\n"
"background:#25383C;}\n")
        self.frame.setFrameShape(QtWidgets.QFrame.StyledPanel)
        self.frame.setFrameShadow(QtWidgets.QFrame.Raised)
        self.frame.setObjectName("frame")
        self.label = QtWidgets.QLabel(self.frame)
        self.label.setGeometry(QtCore.QRect(20, 20, 291, 51))
        self.label.setStyleSheet("font: 63 28pt \"Bahnschrift SemiBold Condensed\";\n"
"color:#56C0EE;")
        self.label.setObjectName("label")
        self.searchquerytextEdit = QtWidgets.QLineEdit(self.frame)
        self.searchquerytextEdit.setGeometry(QtCore.QRect(20, 130, 411, 31))
        self.searchquerytextEdit.setObjectName("searchquerytextEdit")
        self.searchQuerypushButton = QtWidgets.QPushButton(self.frame)
        self.searchQuerypushButton.setGeometry(QtCore.QRect(140, 170, 151, 41))
        self.searchQuerypushButton.setAutoFillBackground(False)
        self.searchQuerypushButton.setStyleSheet("QPushButton\n"
"{\n"
"color:#FFFFFF;\n"
"background:#56C0EE;\n"
"border-radius:20px;\n"
"font: 63 20pt \"Bahnschrift SemiBold Condensed\";\n"
"}"
"QLineEdit\n"
"{\n"
"border-radius:15px;\n"
"color:#56C0EE;\n"
"}")
        self.searchQuerypushButton.setObjectName("searchQuerypushButton")
        self.searchQuerypushButton.clicked.connect(self.search)
        self.showresultslistWidget = QtWidgets.QListWidget(self.frame)
        self.showresultslistWidget.setGeometry(QtCore.QRect(20, 260, 421, 150))
        self.showresultslistWidget.setObjectName("showresultslistWidget")
        self.showresultslistWidget.setStyleSheet("color:white;")

        self.showresultslistWidget.itemDoubleClicked.connect(self.get_tag)
        self.showresultslistWidget.itemClicked.connect(self.get_tag)

        self.label_2 = QtWidgets.QLabel(self.frame)
        self.label_2.setGeometry(QtCore.QRect(20, 220, 291, 31))
        self.label_2.setStyleSheet("font: 63 22pt \"Bahnschrift SemiBold Condensed\";\n"
"color:#56C0EE;")

        self.label_4 = QtWidgets.QLabel(self.frame)
        self.label_4.setGeometry(QtCore.QRect(310, 220, 291, 31))
        self.label_4.setStyleSheet("font: 20 12pt \"Bahnschrift SemiBold Condensed\";\n"
"color:#56C0EE;")



        self.label_2.setObjectName("label_2")
        self.previewlistView = QtWidgets.QPlainTextEdit(self.frame)
        self.previewlistView.setGeometry(QtCore.QRect(20, 420, 320, 200))
        self.previewlistView.setObjectName("previewlistView")
        self.previewlistView.setStyleSheet("color:white;")
        
        
        self.previewpushButton_2 = QtWidgets.QPushButton(self.frame)
        self.previewpushButton_2.setGeometry(QtCore.QRect(350, 500, 100, 41))
        self.previewpushButton_2.setAutoFillBackground(False)
        self.previewpushButton_2.setStyleSheet("QPushButton\n"
"{\n"
"color:#FFFFFF;\n"
"background:#56C0EE;\n"
"border-radius:20px;\n"
"font: 63 15pt \"Bahnschrift SemiBold Condensed\";\n"
"}")
        self.previewpushButton_2.setObjectName("previewpushButton_2")
        self.previewpushButton_2.clicked.connect(self.preview)

        self.summarypushButton_2 = QtWidgets.QPushButton(self.frame)
        self.summarypushButton_2.setGeometry(QtCore.QRect(350, 440, 100, 41))
        self.summarypushButton_2.setAutoFillBackground(False)
        self.summarypushButton_2.setStyleSheet("QPushButton\n"
"{\n"
"color:#FFFFFF;\n"
"background:#56C0EE;\n"
"border-radius:20px;\n"
"font: 63 15pt \"Bahnschrift SemiBold Condensed\";\n"
"}")
        self.summarypushButton_2.setObjectName("summarypushButton_2")
        self.summarypushButton_2.clicked.connect(self.Summary)

        self.comboBox = QtWidgets.QComboBox(self.frame)
        self.comboBox.setGeometry(QtCore.QRect(350, 224, 90, 25))
        self.comboBox.setStyleSheet("QComboBox\n"
"{\n"
"color:#FFFFFF;\n"
"background:#56C0EE;\n"
"border-radius:20px;\n"
"font: 12 12pt \"Bahnschrift SemiBold Condensed\";\n"
"}")
        self.comboBox.setObjectName("comboBox")
        self.comboBox.addItem("Technology")
        self.comboBox.addItem("Arts")
        self.comboBox.addItem("Biology")
        self.comboBox.addItem("Chemistry")
        self.comboBox.addItem("Economy")
        self.comboBox.addItem("Literature")
        self.comboBox.addItem("Physics")
        self.comboBox.addItem("Sports")
        self.comboBox.addItem("Others")
        self.comboBox.currentTextChanged.connect(self.on_combobox_changed)






        
        self.openfilepushButton_3 = QtWidgets.QPushButton(self.frame)
        self.openfilepushButton_3.setGeometry(QtCore.QRect(350, 560, 100, 41))
        self.openfilepushButton_3.setAutoFillBackground(False)
        self.openfilepushButton_3.setStyleSheet("QPushButton\n"
"{\n"
"color:#FFFFFF;\n"
"background:#56C0EE;\n"
"border-radius:20px;\n"
"font: 63 15pt \"Bahnschrift SemiBold Condensed\";\n"
"}")
        self.openfilepushButton_3.setObjectName("openfilepushButton_3")
        self.openfilepushButton_3.clicked.connect(self.open_a_file)
        self.label_3 = QtWidgets.QLabel(self.frame)
        self.label_3.setGeometry(QtCore.QRect(350, 20, 81, 101))
        self.label_3.setText("")
        self.label_3.setPixmap(QtGui.QPixmap("GUI_img/doc_icon.png"))
        self.label_3.setScaledContents(True)
        self.label_3.setObjectName("label_3")
        MainWindow.setCentralWidget(self.centralwidget)
        self.menubar = QtWidgets.QMenuBar(MainWindow)
        self.menubar.setGeometry(QtCore.QRect(0, 0, 457, 21))
        self.menubar.setObjectName("menubar")
        MainWindow.setMenuBar(self.menubar)
        self.statusbar = QtWidgets.QStatusBar(MainWindow)
        self.statusbar.setObjectName("statusbar")
        MainWindow.setStatusBar(self.statusbar)


        self.retranslateUi(MainWindow)
        QtCore.QMetaObject.connectSlotsByName(MainWindow)

    def retranslateUi(self, MainWindow):
        _translate = QtCore.QCoreApplication.translate
        MainWindow.setWindowTitle(_translate("MainWindow", "MainWindow"))
        self.label.setText(_translate("MainWindow", "Search Files"))
        self.searchQuerypushButton.setText(_translate("MainWindow", "SEARCH"))
        self.label_2.setText(_translate("MainWindow", "MATCH FOUND"))
        self.label_4.setText(_translate("MainWindow", "TAG:"))
        self.previewpushButton_2.setText(_translate("MainWindow", "Preview"))
        self.summarypushButton_2.setText(_translate("MainWindow", "Summary"))
        self.openfilepushButton_3.setText(_translate("MainWindow", "Open File"))
class Ui_MainWindow(object):
    
    def search_window(self):
        #print('search some')
        self.window = QtWidgets.QMainWindow()
        self.ui = Ui_SearchWindow()
        self.ui.setupUi(self.window)
        self.window.show()
        MainWindow.hide()
        
        
    
    def scanner(self):
        with open('status/path.txt','r') as f:
                self.datadir = str(f.read())
        f.close()
        self.ui.progressBar.setValue(0)
        self.count=0
        
        ##########################################################################
        value = self.ui.progressBar.value()
        
        datadir = self.datadir
        self.status_directory = 'status/' 

        if not os.path.isdir(self.status_directory):
            #print('make status dir')
            os.mkdir('status')


        first_time = False
        self.search_table = {'file_name':[],'keys':[],'hash_map':[],'text':[],'tag':[],'location':[]}
        self.table = pd.DataFrame(self.search_table)

        self.all_files =  []
        for each_dir in self.datadir.split(' | '):
            for root,folder,files in os.walk(each_dir):
                for file in files:
                    if file.split('.')[-1] in ['pdf','pptx','docx','txt','csv','xlsx']:
                 	    self.all_files.append(root+'/'+file)

        if "done_list.csv" not in os.listdir(self.status_directory):
            #print('file created ')
            pd.DataFrame({'files':[]}).to_csv(self.status_directory+'done_list.csv',index=False)
            self.table.to_excel(self.status_directory+'table.xlsx')
            self.done_list = pd.read_csv(self.status_directory+'done_list.csv')

        else:
            self.done_list = pd.read_csv(self.status_directory+'done_list.csv')
            self.table = pd.read_excel(self.status_directory+'table.xlsx',index=False)
        #table.set_index('file_name')

        self.files_to_process = list(set(self.all_files)-set(self.done_list.files))

        start_ini = time.time()
        #print(self.files_to_process)
        self.total_length = len(self.files_to_process)
        
        
        
        self.timer = QTimer()        
        self.timer.timeout.connect(self.handleTimer)
        self.timer.start(100)

    def handleTimer(self):
        #print("called")
        value = self.ui.progressBar.value()
            #try:
        if self.count>self.total_length-1:
            self.count-=1
        if len(self.files_to_process)==0:
            self.timer.stop()
            self.search_window()
            return 
        file = self.files_to_process[self.count]
        self.count+=1
        
        if True:
                if file.endswith(".pdf"):
                    start = time.time()
                    text = convert_pdf_to_txt(file)
                    text = text.replace("\n"," ")
                    text = text.replace("\\x"," ")
                    self.search_table['text'].append(text)
                    self.search_table['tag'].append(tag(text))
                    self.search_table['location'].append(file)

                    text = " ".join(re.split("[^0-9a-zA-Z]",text.lower() ))
                    
                    text = text.split(" ")
                    if len(text)<100:
                        self.done_list = self.done_list.append({'files':file},ignore_index = True)
                        #print('bad QUALITY can`t read:',file)
                        self.ui.listWidget.addItem('bad QUALITY can`t read: '+file.split('/')[-1])
                        
                        #continue
                    
                    count = Counter(text)

                    #print(file," scaned in ",time.time()-start ,"sec")
                    self.ui.listWidget.addItem(file.split('/')[-1]+" scanned in "+str(np.round(time.time()-start,decimals=2)) +"sec")

                    hash_file = hash_table(count) 



                    self.search_table['file_name'].append(file.split('/')[-1])
                    self.search_table['keys'].append(list(hash_file.keys()))
                    self.search_table['hash_map'].append(hash_file)
                    self.done_list = self.done_list.append({'files':file},ignore_index = True)

                if file.endswith(".docx"):
                    start = time.time()
                    text = docx_read(file)
                    text = text.replace("\n"," ")
                    self.search_table['text'].append(text)
                    self.search_table['tag'].append(tag(text))
                    self.search_table['location'].append(file)

                    text = " ".join(re.split("[^0-9a-zA-Z]",text.lower() ))
                    
                    text = text.split(" ")
                    if len(text)<100:
                        #print('bad QUALITY can`t read:',file)
                        self.done_list = self.done_list.append({'files':file},ignore_index = True)
                        self.ui.listWidget.addItem('bad QUALITY can`t read: '+file.split('/')[-1])
                        #continue
                   
                    count = Counter(text)

                    #print(file," scaned in ",time.time()-start ,"sec")
                    self.ui.listWidget.addItem(file.split('/')[-1]+" scanned in "+str(np.round(time.time()-start,decimals=2))+"sec")

                    hash_file = hash_table(count) 

                    self.search_table['file_name'].append(file.split('/')[-1])
                    self.search_table['keys'].append(list(hash_file.keys()))
                    self.search_table['hash_map'].append(hash_file)
                    self.done_list = self.done_list.append({'files':file},ignore_index = True)

                if file.endswith(".txt"):
                    start = time.time()
                    text = txt_read(file)
                    text = text.replace("\n"," ")
                    self.search_table['text'].append(text)
                    self.search_table['tag'].append(tag(text))
                    self.search_table['location'].append(file)

                    text = " ".join(re.split("[^0-9a-zA-Z]",text.lower() ))
                    
                    text = text.split(" ")
                    if len(text)<100:
                        #print('bad QUALITY can`t read:',file)
                        self.done_list = self.done_list.append({'files':file},ignore_index = True)
                        self.ui.listWidget.addItem('bad QUALITY can`t read: '+file.split('/')[-1])
                        #continue
                   
                    count = Counter(text)

                    #print(file," scaned in ",time.time()-start ,"sec")
                    self.ui.listWidget.addItem(file.split('/')[-1]+" scanned in "+str(np.round(time.time()-start,decimals=2))+"sec")

                    hash_file = hash_table(count) 

                    self.search_table['file_name'].append(file.split('/')[-1])
                    self.search_table['keys'].append(list(hash_file.keys()))
                    self.search_table['hash_map'].append(hash_file)
                    self.done_list = self.done_list.append({'files':file},ignore_index = True)

                if file.endswith(".csv"):
                    start = time.time()
                    text = csv_read(file)
                    text = text.replace("\n"," ")

                    
                    self.search_table['tag'].append('Others')
                    text = " ".join(re.split("[^a-zA-Z]",text.lower() ))
                    self.search_table['text'].append(text)
                    self.search_table['location'].append(file)

                    text = text.split(" ")
                    if len(text)<100:
                        #print('bad QUALITY can`t read:',file)
                        self.done_list = self.done_list.append({'files':file},ignore_index = True)
                        self.ui.listWidget.addItem('bad QUALITY can`t read: '+file.split('/')[-1])
                        #continue
                   
                    count = Counter(text)

                    #print(file," scaned in ",time.time()-start ,"sec")
                    self.ui.listWidget.addItem(file.split('/')[-1]+" scanned in "+str(np.round(time.time()-start,decimals=2))+"sec")

                    hash_file = hash_table(count) 

                    self.search_table['file_name'].append(file.split('/')[-1])
                    self.search_table['keys'].append(list(hash_file.keys()))
                    self.search_table['hash_map'].append(hash_file)
                    self.done_list = self.done_list.append({'files':file},ignore_index = True)
                if file.endswith(".xlsx"):
                    start = time.time()
                    text = excel_read(file)
                    text = text.replace("\n"," ")
                    
                    self.search_table['tag'].append('Others')
                    text = " ".join(re.split("[^a-zA-Z]",text.lower() ))
                    self.search_table['text'].append(text)
                    self.search_table['location'].append(file)

                    text = text.split(" ")
                    if len(text)<100:
                        #print('bad QUALITY can`t read:',file)
                        self.done_list = self.done_list.append({'files':file},ignore_index = True)
                        self.ui.listWidget.addItem('bad QUALITY can`t read: '+file.split('/')[-1])
                        #continue
                   
                    count = Counter(text)

                    #print(file," scaned in ",time.time()-start ,"sec")
                    self.ui.listWidget.addItem(file.split('/')[-1]+" scanned in "+str(np.round(time.time()-start,decimals=2))+"sec")

                    hash_file = hash_table(count) 

                    self.search_table['file_name'].append(file.split('/')[-1])
                    self.search_table['keys'].append(list(hash_file.keys()))
                    self.search_table['hash_map'].append(hash_file)
                    self.done_list = self.done_list.append({'files':file},ignore_index = True)

                if file.endswith(".pptx"):
                    start = time.time()
                    text = ppt_read(file)
                    
                    text = text.replace("\n"," ")
                    self.search_table['text'].append(text)
                    self.search_table['tag'].append(tag(text))
                    self.search_table['location'].append(file)


                    text = " ".join(re.split("[^0-9a-zA-Z]",text.lower() ))
                    
                    text = text.split(" ")
                    if len(text)<100:
                        self.done_list = self.done_list.append({'files':file},ignore_index = True)
                        #print('bad QUALITY can`t read:',file)
                        self.ui.listWidget.addItem('bad QUALITY can`t read: '+file.split('/')[-1])
                        #continue
                    
                    count = Counter(text)

                    #print(file," scaned in ",time.time()-start ,"sec")
                    self.ui.listWidget.addItem(file.split('/')[-1]+" scanned in "+str(np.round(time.time()-start,decimals=2)) +"sec")

                    hash_file = hash_table(count) 

                    self.search_table['file_name'].append(file.split('/')[-1])
                    self.search_table['keys'].append(list(hash_file.keys()))
                    self.search_table['hash_map'].append(hash_file)
                    self.done_list = self.done_list.append({'files':file},ignore_index = True)
                else:
                    self.done_list = self.done_list.append({'files':file},ignore_index = True)
        
    
        if value < 100:
            value = self.count*100/self.total_length
            self.ui.progressBar.setValue(value)
        else:
            pre_table = pd.read_excel(self.status_directory+'table.xlsx',index=False)
            self.table = pd.DataFrame(self.search_table)
            self.table = pre_table.append(self.table)

            self.table = self.table.reset_index()
            self.table = self.table[['file_name','keys','hash_map','text','tag','location']]

            for i in self.table.index:
                if type(self.table.loc[i,'keys'])==str:
                    #print(self.table.at[i,'file_name'])
                    self.table.at[i,'keys'] = eval(self.table.loc[i,'keys'])
                    try:
                    	self.table.at[i,'hash_map'] = eval(self.table.loc[i,'hash_map'])
                    except:
                    	pass

            self.table.to_excel(self.status_directory+'table.xlsx')
            self.done_list= self.done_list[['files']]
            
            self.done_list.to_csv(self.status_directory+'done_list.csv')
            
            self.timer.stop()
            self.search_window()
            
    
    def scan_file(self):
        self.t = 0
        #print("clicked")
        #print(self.lineEdit.text())
        temp_text  = self.lineEdit.text()
        for i,path in enumerate(temp_text.split(' | ')):
            if not  os.path.isdir(path):
                 self.lineEdit.setText(f"{i+1} DOCUMENTS PATH IS WRONG ENTER CORRECT PATH")
                 return False

        with open('status/path.txt','w') as f:
            f.write(str(temp_text))
        f.close()
        self.window = QtWidgets.QMainWindow()
        self.ui = Ui_ScanWindow()
        self.ui.setupUi(self.window)
        self.window.show()
        MainWindow.hide()
        self.scanner()
	    	

        
    def setupUi(self, MainWindow):
        MainWindow.setObjectName("MainWindow")
        MainWindow.resize(458, 643)

        self.centralwidget = QtWidgets.QWidget(MainWindow)
        self.centralwidget.setObjectName("centralwidget")
        self.frame = QtWidgets.QFrame(self.centralwidget)
        self.frame.setGeometry(QtCore.QRect(-10, -30, 471, 671))

        self.frame.setStyleSheet("QFrame\n"
"{\n"
"color:#54B0EE;\n"
"background:#25383C;\n"
"}\n"
"QProgressBar\n"
"{\n"
"color:#54B0EE;\n"
"background:54B0EE;\n"
"}\n"
                               
)
        self.frame.setFrameShape(QtWidgets.QFrame.StyledPanel)
        self.frame.setFrameShadow(QtWidgets.QFrame.Raised)
        self.frame.setObjectName("frame")
        self.label = QtWidgets.QLabel(self.frame)
        self.label.setGeometry(QtCore.QRect(50, 80, 281, 71))
        self.label.setAutoFillBackground(False)
        self.label.setStyleSheet("\n"
"\n"
"font: 63 28pt \"Bahnschrift SemiBold Condensed\";")
        self.label.setObjectName("label")
        self.label_2 = QtWidgets.QLabel(self.frame)
        self.label_2.setGeometry(QtCore.QRect(50, 160, 131, 161))
        self.label_2.setText("")
        self.label_2.setPixmap(QtGui.QPixmap("GUI_img/doc_icon.png"))
        self.label_2.setScaledContents(True)
        self.label_2.setObjectName("label_2")
        self.label_3 = QtWidgets.QLabel(self.frame)
        self.label_3.setGeometry(QtCore.QRect(50, 350, 121, 21))
        self.label_3.setStyleSheet("\n"
"font: 63 12pt \"Bahnschrift SemiBold Condensed\";")
        self.label_3.setObjectName("label_3")
        self.pushButton = QtWidgets.QPushButton(self.frame)
        self.pushButton.setGeometry(QtCore.QRect(50, 420, 151, 41))
        self.pushButton.setAutoFillBackground(False)
        self.pushButton.setStyleSheet("QPushButton\n"
"{\n"
"color:#FFFFFF;\n"
"background:#56C0EE;\n"
"border-radius:20px;\n"
"font: 63 17pt \"Bahnschrift SemiBold Condensed\";\n"
"}")
        self.pushButton.setObjectName("pushButton")
        self.pushButton.clicked.connect(self.scan_file)
        self.lineEdit = QtWidgets.QLineEdit(self.frame)
        self.lineEdit.setGeometry(QtCore.QRect(50, 380, 371, 31))
        self.lineEdit.setObjectName("lineEdit")
        self.lineEdit.setStyleSheet("QLineEdit\n"
"{\n"
"border-radius:15px;\n"
"color:#56C0EE;\n" 
                        
"}" )
        self.status_directory = 'status/' 

        if not os.path.isdir(self.status_directory):
            #print('make status dir')
            os.mkdir('status')
        if 'path.txt' in os.listdir('status'):
            with open('status/path.txt','r') as f:
                text  = str(f.read())
                self.lineEdit.setText(text)
            f.close()
        else:
            open("status/path.txt", "w")
        
        MainWindow.setCentralWidget(self.centralwidget)
        self.menubar = QtWidgets.QMenuBar(MainWindow)
        self.menubar.setGeometry(QtCore.QRect(0, 0, 458, 21))
        self.menubar.setObjectName("menubar")
        MainWindow.setMenuBar(self.menubar)
        self.statusbar = QtWidgets.QStatusBar(MainWindow)
        self.statusbar.setObjectName("statusbar")
        MainWindow.setStatusBar(self.statusbar)

        self.retranslateUi(MainWindow)
        QtCore.QMetaObject.connectSlotsByName(MainWindow)

    def retranslateUi(self, MainWindow):
        _translate = QtCore.QCoreApplication.translate
        MainWindow.setWindowTitle(_translate("MainWindow", "MainWindow"))
        self.label.setText(_translate("MainWindow", "FILE BUDDY"))
        self.label_3.setText(_translate("MainWindow", "FOLDER NAME"))
        self.pushButton.setText(_translate("MainWindow", "START SCAN"))


class Ui_ScanWindow(object):
    def setupUi(self, MainWindow):
        MainWindow.setObjectName("MainWindow")
        MainWindow.resize(457, 646)
        self.centralwidget = QtWidgets.QWidget(MainWindow)
        self.centralwidget.setObjectName("centralwidget")
        self.frame = QtWidgets.QFrame(self.centralwidget)
        self.frame.setGeometry(QtCore.QRect(0, 0, 481, 631))
        self.frame.setStyleSheet("QFrame\n"
"{\n"
"background:#25383C;\n"
"}")
        self.frame.setFrameShape(QtWidgets.QFrame.StyledPanel)
        self.frame.setFrameShadow(QtWidgets.QFrame.Raised)
        self.frame.setObjectName("frame")
        self.label = QtWidgets.QLabel(self.frame)
        self.label.setGeometry(QtCore.QRect(40, 130, 411, 91))
        self.label.setStyleSheet("font: 63 28pt \"Bahnschrift SemiBold Condensed\";\n"
"color:#56C0EE;")
        self.label.setObjectName("label")
        self.progressBar = QtWidgets.QProgressBar(self.frame)
        self.progressBar.setGeometry(QtCore.QRect(40, 240, 381, 23))
        self.progressBar.setProperty("value", 24)
        self.progressBar.setObjectName("progressBar")
        self.listWidget = QtWidgets.QListWidget(self.frame)
        self.listWidget.setGeometry(QtCore.QRect(40, 290, 351, 241))
        self.listWidget.setObjectName("listWidget")
        self.listWidget.setStyleSheet("color:white;")
        MainWindow.setCentralWidget(self.centralwidget)
        self.menubar = QtWidgets.QMenuBar(MainWindow)
        self.menubar.setGeometry(QtCore.QRect(0, 0, 457, 21))
        self.menubar.setObjectName("menubar")
        MainWindow.setMenuBar(self.menubar)
        self.statusbar = QtWidgets.QStatusBar(MainWindow)
        self.statusbar.setObjectName("statusbar")
        MainWindow.setStatusBar(self.statusbar)

        self.retranslateUi(MainWindow)
        QtCore.QMetaObject.connectSlotsByName(MainWindow)

    def retranslateUi(self, MainWindow):
        _translate = QtCore.QCoreApplication.translate
        MainWindow.setWindowTitle(_translate("MainWindow", "MainWindow"))
        self.label.setText(_translate("MainWindow", "SCANNING FILES"))


# In[2]:

cv = pickle.load(open('save_model/cv.pickel','rb'))
tagger = pickle.load(open('save_model/tagger.pickel','rb'))

def tag(text):
    try:
        return tagger.predict(cv.transform([text]).toarray())[0]
    except:
        return 'Others'

def collect(data):
    s = []
    for i in data:
        s.append(list(i.keys())[0])
    return s
def order_calc(search_order,present_order):
    if len(search_order)==1:
        return 1
    else:
        s = {key:i for i,key in enumerate(search_order)}
        p=dict()
        count=0
        for i in range(len(present_order)-1):
            j = s[present_order[i]]
            if j<len(search_order)-1:
                if present_order[i+1]==search_order[j+1]:
                    count+=1
        return count
def spacer(string):
    text = ''
    if len(string)>0:
        for i,val in enumerate(string[:-1]):
            if string[i].islower() and string[i+1].isupper():
                text=text + string[i]+' '
            else:
                text = text + string[i]
        return text+string[-1]
    else:
        return ""
def ppt_read(path):
    text = ''
    prs = Presentation(path)
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                text  = text + shape.text
            except:
                pass
    return text
def docx_read(path):
    f = open(path, 'rb')
    document = Document(f)
    f.close()
    text = ''
    for lines in document.paragraphs:
        text = text+lines.text
    return text

def csv_read(path):
     f = pd.read_csv(path)    
     text = f.to_string()
     return text
def excel_read(path):
     f = pd.read_excel(path)    
     text = f.to_string()
     return text


def txt_read(path):
    f = open(path, 'rb')
    document = f.read()
    document = str(document)
    f.close()
    return document 

def convert_pdf_to_txt(path):
    rsrcmgr = PDFResourceManager()
    retstr = io.StringIO()
    codec = 'utf-8'
    laparams = LAParams()
    device = TextConverter(rsrcmgr, retstr, laparams=laparams)
    fp = open(path, 'rb')
    interpreter = PDFPageInterpreter(rsrcmgr, device)
    password = ""
    maxpages = 5
    caching = True
    pagenos = set()
    try:

	    for page in PDFPage.get_pages(fp, pagenos, maxpages=maxpages,
	                                  password=password,
	                                  caching=caching,
	                                  check_extractable=True):
	        interpreter.process_page(page)
    except:

	    retstr.getvalue()

    text = retstr.getvalue()

    fp.close()
    device.close()
    retstr.close()
    return text
def hash_table(text):
    hash_map = dict()
    for i,val in enumerate(text):
        l = len(val)
        if l>2:
            key = query(val)
            if key in hash_map.keys():
                hash_map[key][val]=text[val]
            else:
                hash_map[key]=dict({val:text[val]})
                
    return hash_map
def query(x):
    if len(x)>0:
        return sum(map(ord,x.lower()))/len(x)



if __name__ == "__main__":
    import sys
    app = QtWidgets.QApplication(sys.argv)
    app.setWindowIcon(QtGui.QIcon("GUI_img/doc_icon.png"))
    MainWindow = QtWidgets.QMainWindow()
    ui = Ui_MainWindow()
    ui.setupUi(MainWindow)
    MainWindow.show()
    sys.exit(app.exec_())



